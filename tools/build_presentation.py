"""
Build a 16-slide class-project presentation summarising the work.

Layout:
  Slide 1     – title
  Slide 2     – Problem (left = architecture placeholder, right = text)
  Slide 3-14  – split layout: explanation/analysis on left, figure on right
  Slide 15    – Conclusion (text-only)
  Slide 16    – Future Work (text-only)

Outputs:
  PAPER_PLOTS/figures/class_distribution.pdf/png
  PAPER_PLOTS/presentation/AV_VLM_Benchmark.pptx
"""

import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT       = Path(__file__).resolve().parent.parent
FIG_DIR    = ROOT / "PAPER_PLOTS" / "main"
APP_FIG    = ROOT / "PAPER_PLOTS" / "appendix"
PRES_DIR   = ROOT / "PAPER_PLOTS" / "presentation"
PRES_DIR.mkdir(parents=True, exist_ok=True)

DS_JSON = ROOT / "Data" / "dataset.json"

# ── Configuration ─────────────────────────────────────────────────────────────
TITLE         = "Benchmarking Vision-Language Models\nfor Anomaly Detection in Autonomous Driving"
SUBTITLE      = "An empirical study of 9 VLMs and 6 baselines\non a 15,000-image dashcam benchmark"
AUTHOR        = "Mahyar Ghazanfari"
AFFILIATION   = "[Your Affiliation]"
DATE_STR      = "April 2026"

# Colour palette (kept simple and modern)
NAVY    = RGBColor(0x1A, 0x3B, 0x5C)   # primary
ACCENT  = RGBColor(0xE6, 0x4A, 0x19)   # accent (highlights)
GREY_D  = RGBColor(0x33, 0x33, 0x33)   # dark text
GREY_M  = RGBColor(0x66, 0x66, 0x66)   # secondary text
GREY_L  = RGBColor(0xCC, 0xCC, 0xCC)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

FONT_FAM = "Calibri"   # broadly-supported clean sans-serif

TOTAL_SLIDES = 16

# Slide dimensions (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ──────────────────────────────────────────────────────────────────────────────
# 1. Generate dataset class-distribution figure (we don't have one as PDF)
# ──────────────────────────────────────────────────────────────────────────────
print("── Generating dataset_class_distribution figure ──")

ANOMALY_CLASSES = [
    "animal_on_road", "extreme_weather", "road_surface_hazard",
    "fallen_debris_or_vegetation", "strange_object_on_road",
    "vehicle_incident", "infrastructure_failure", "human_presence_anomaly",
    "adverse_lighting", "oversized_or_unusual_vehicle", "multi_hazard_compound",
]
CLASS_DISPLAY = {
    "animal_on_road":              "Animal on Road",
    "extreme_weather":             "Extreme Weather",
    "road_surface_hazard":         "Road Surface Hazard",
    "fallen_debris_or_vegetation": "Fallen Debris/Veg.",
    "strange_object_on_road":      "Strange Object",
    "vehicle_incident":            "Vehicle Incident",
    "infrastructure_failure":      "Infrastructure Fail.",
    "human_presence_anomaly":      "Human Presence",
    "adverse_lighting":            "Adverse Lighting",
    "oversized_or_unusual_vehicle":"Oversized Vehicle",
    "multi_hazard_compound":       "Multi-Hazard",
}

with open(DS_JSON) as f:
    ds = json.load(f)
samples = ds.get("samples", ds)
class_counts  = {c: 0 for c in ANOMALY_CLASSES}
n_anomalous   = 0
for r in samples.values():
    if isinstance(r, dict) and r.get("anomaly_present"):
        n_anomalous += 1
        ac = r.get("anomaly_class")
        if ac in class_counts:
            class_counts[ac] += 1
n_normal = sum(1 for r in samples.values()
               if isinstance(r, dict) and not r.get("anomaly_present"))

plt.rcParams.update({
    "figure.facecolor": "white",
    "axes.facecolor":   "white",
    "font.family":      "sans-serif",
    "axes.spines.top":  False,
    "axes.spines.right":False,
})

# Two-panel: left = normal vs anomalous, right = per-class bars
fig, axes = plt.subplots(1, 2, figsize=(12, 5),
                          gridspec_kw={"width_ratios": [1, 2.4]})

# Panel A: normal vs anomalous
ax = axes[0]
counts  = [n_normal, n_anomalous]
labels  = [f"Normal\n({n_normal:,})", f"Anomalous\n({n_anomalous:,})"]
colors  = ["#4DA8DA", "#E64A19"]
bars = ax.bar(labels, counts, color=colors, width=0.6, edgecolor="white", linewidth=2)
for bar, c in zip(bars, counts):
    ax.text(bar.get_x() + bar.get_width() / 2, c + 200,
            f"{c:,}", ha="center", fontsize=11, fontweight="bold")
ax.set_ylabel("Number of images", fontsize=11)
ax.set_title("Image-level split", fontsize=12, pad=10)
ax.set_ylim(0, max(counts) * 1.15)

# Panel B: per-class
ax = axes[1]
order = sorted(ANOMALY_CLASSES, key=lambda c: -class_counts[c])
xs    = [CLASS_DISPLAY[c] for c in order]
ys    = [class_counts[c] for c in order]
class_colors = plt.cm.viridis_r(range(len(order)))
bars = ax.barh(xs, ys, color=class_colors, edgecolor="white", linewidth=0.5)
for bar, v in zip(bars, ys):
    ax.text(v + max(ys) * 0.01, bar.get_y() + bar.get_height() / 2,
            f"{v:,}", va="center", fontsize=9)
ax.set_xlabel("Number of anomalous images", fontsize=11)
ax.set_title("Anomaly class distribution (11 classes)", fontsize=12, pad=10)
ax.invert_yaxis()
plt.tight_layout()

dist_path_pdf = ROOT / "PAPER_PLOTS" / "main" / "fig00_dataset_distribution.pdf"
dist_path_png = ROOT / "PAPER_PLOTS" / "main" / "fig00_dataset_distribution.png"
fig.savefig(dist_path_pdf, dpi=300, bbox_inches="tight")
fig.savefig(dist_path_png, dpi=200, bbox_inches="tight")
plt.close(fig)
print(f"  Saved: {dist_path_pdf.name}")


# ──────────────────────────────────────────────────────────────────────────────
# 2. Build the presentation
# ──────────────────────────────────────────────────────────────────────────────
print("\n── Building presentation ──")

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank layout we lay out manually


def set_text(tf, text, size=14, bold=False, color=GREY_D, align=PP_ALIGN.LEFT,
             font=FONT_FAM, line_spacing=1.2):
    """Write a single paragraph into a textframe."""
    tf.word_wrap = True
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_paragraph(tf, text, size=14, bold=False, color=GREY_D,
                  align=PP_ALIGN.LEFT, font=FONT_FAM, space_before=0,
                  bullet=False, line_spacing=1.25):
    """Append a paragraph to a text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line_spacing
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    if bullet:
        run.text = "•  " + text
    else:
        run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_box(slide, left, top, width, height, fill=None, line=None):
    from pptx.shapes.autoshape import Shape
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
    box.shadow.inherit = False
    return box


# ──────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title slide
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)

# Background accent strip on top
add_box(slide, 0, 0, SLIDE_W, Inches(0.4), fill=NAVY)
add_box(slide, 0, SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.25), fill=ACCENT)

# Title
title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0),
                                      SLIDE_W - Inches(1.6), Inches(2.0))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.line_spacing = 1.15
for line in TITLE.split("\n"):
    if p.text == "" and not p.runs:
        run = p.add_run()
    else:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.15
        run = p.add_run()
    run.text = line
    run.font.name = FONT_FAM
    run.font.size = Pt(38)
    run.font.bold = True
    run.font.color.rgb = NAVY

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.3),
                                    SLIDE_W - Inches(1.6), Inches(1.0))
sub_tf = sub_box.text_frame
sub_tf.word_wrap = True
sub_p = sub_tf.paragraphs[0]
sub_p.alignment = PP_ALIGN.CENTER
sub_p.line_spacing = 1.3
first = True
for line in SUBTITLE.split("\n"):
    if first:
        run = sub_p.add_run()
        first = False
    else:
        sub_p = sub_tf.add_paragraph()
        sub_p.alignment = PP_ALIGN.CENTER
        sub_p.line_spacing = 1.3
        run = sub_p.add_run()
    run.text = line
    run.font.name = FONT_FAM
    run.font.size = Pt(20)
    run.font.italic = True
    run.font.color.rgb = GREY_M

# Author + affiliation
auth_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.7),
                                     SLIDE_W - Inches(1.6), Inches(1.2))
atf = auth_box.text_frame
set_text(atf, AUTHOR, size=24, bold=True, color=GREY_D, align=PP_ALIGN.CENTER)
add_paragraph(atf, AFFILIATION, size=18, color=GREY_M,
              align=PP_ALIGN.CENTER, space_before=4)
add_paragraph(atf, DATE_STR, size=15, color=GREY_M,
              align=PP_ALIGN.CENTER, space_before=8)


# ──────────────────────────────────────────────────────────────────────────────
# Helper to add a content slide (text left, figure right)
# ──────────────────────────────────────────────────────────────────────────────
def add_content_slide(slide_num, title, subtitle, bullets, figure_path,
                      figure_caption=None):
    """
    slide_num     : int (for footer)
    title         : str — main slide title
    subtitle      : str — small subtitle below title (or "" to skip)
    bullets       : list[str] — bullet points (analysis / takeaways)
    figure_path   : Path to PDF or PNG figure
    figure_caption: optional small caption under the figure
    """
    slide = prs.slides.add_slide(BLANK)

    # Top stripe
    add_box(slide, 0, 0, SLIDE_W, Inches(0.18), fill=NAVY)

    # Title
    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.30),
                                      SLIDE_W - Inches(1.0), Inches(0.7))
    set_text(t_box.text_frame, title, size=26, bold=True, color=NAVY,
             align=PP_ALIGN.LEFT)

    # Subtitle (optional)
    if subtitle:
        s_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0),
                                         SLIDE_W - Inches(1.0), Inches(0.4))
        set_text(s_box.text_frame, subtitle, size=14, color=GREY_M, align=PP_ALIGN.LEFT)

    # Left column: bullets / explanation
    LEFT_X     = Inches(0.55)
    LEFT_Y     = Inches(1.6) if subtitle else Inches(1.3)
    LEFT_W     = Inches(5.6)
    LEFT_H     = SLIDE_H - LEFT_Y - Inches(0.6)

    text_box = slide.shapes.add_textbox(LEFT_X, LEFT_Y, LEFT_W, LEFT_H)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)

    # First bullet uses paragraphs[0]
    first = True
    for b in bullets:
        if first:
            set_text(tf, "•  " + b, size=15, color=GREY_D,
                     line_spacing=1.35)
            first = False
        else:
            add_paragraph(tf, b, size=15, color=GREY_D, bullet=True,
                          space_before=8, line_spacing=1.35)

    # Right column: figure
    RIGHT_X = Inches(6.4)
    RIGHT_Y = Inches(1.3)
    RIGHT_W = SLIDE_W - RIGHT_X - Inches(0.5)
    RIGHT_H = SLIDE_H - RIGHT_Y - Inches(0.7)

    # python-pptx works best with PNG; prefer .png if available
    fig_path = Path(figure_path)
    if fig_path.suffix == ".pdf":
        png_alt = fig_path.with_suffix(".png")
        if png_alt.exists():
            fig_path = png_alt
    if fig_path.exists():
        pic = slide.shapes.add_picture(str(fig_path),
                                        RIGHT_X, RIGHT_Y,
                                        width=RIGHT_W)
        # If picture is too tall, scale by height instead
        if pic.height > RIGHT_H:
            slide.shapes._spTree.remove(pic._element)
            pic = slide.shapes.add_picture(str(fig_path),
                                            RIGHT_X, RIGHT_Y,
                                            height=RIGHT_H)
            # Re-centre horizontally
            pic.left = RIGHT_X + (RIGHT_W - pic.width) // 2
    else:
        ph = slide.shapes.add_textbox(RIGHT_X, RIGHT_Y, RIGHT_W, RIGHT_H)
        set_text(ph.text_frame, f"[Missing figure: {fig_path.name}]",
                 size=14, color=ACCENT, align=PP_ALIGN.CENTER)

    # Figure caption
    if figure_caption:
        cap_y = SLIDE_H - Inches(0.55)
        cap_box = slide.shapes.add_textbox(RIGHT_X, cap_y,
                                            RIGHT_W, Inches(0.35))
        set_text(cap_box.text_frame, figure_caption, size=10,
                 color=GREY_M, align=PP_ALIGN.CENTER)

    # Footer (slide number)
    footer = slide.shapes.add_textbox(SLIDE_W - Inches(1.0),
                                       SLIDE_H - Inches(0.4),
                                       Inches(0.8), Inches(0.3))
    set_text(footer.text_frame, f"{slide_num} / {TOTAL_SLIDES}",
             size=10, color=GREY_M, align=PP_ALIGN.RIGHT)

    return slide


# ──────────────────────────────────────────────────────────────────────────────
# Slides 2-13: Content slides
# ──────────────────────────────────────────────────────────────────────────────
SLIDES = [
    # (title, subtitle, bullets, figure_path, caption)

    # --- Slide 2: Dataset overview ---
    (
        "Dataset Overview",
        "15,000 dashcam images for autonomous-driving anomaly benchmarking",
        [
            "We curated a balanced benchmark of 15,000 dashcam images: "
            "10,000 normal driving scenes and 5,000 anomalous scenes.",
            "Each anomalous image is annotated with one of 11 hazard "
            "classes plus a free-form description, severity rating, "
            "and recommended driving action.",
            "Class distribution is intentionally long-tailed — "
            "Oversized Vehicle (n=2,037) and Human Presence (n=728) "
            "are the largest classes; Multi-Hazard (n=33) is the rarest.",
            "This long-tail makes per-class evaluation realistic and "
            "exposes weaknesses on rare but safety-critical anomalies.",
        ],
        FIG_DIR / "fig00_dataset_distribution.pdf",
        "Image split (left) and per-class anomaly distribution (right).",
    ),

    # --- Slide 3: Binary detection ---
    (
        "Binary Anomaly Detection",
        "Can VLMs decide whether ANY anomaly is present? (15,000 images)",
        [
            "We evaluate 9 VLMs (zero-shot prompting) plus 6 baselines "
            "(3 CLIP and 3 SigLIP variants) and 1 reconstruction "
            "autoencoder on the binary anomaly-detection task.",
            "Qwen2.5-VL-3B is the clear winner with AUPRC = 0.82 — "
            "even outperforming much larger models like InternVL3-8B "
            "and LLaMA-3.2-11B.",
            "Most VLMs over-predict 'anomaly' (precision suffers); "
            "the autoencoder underperforms badly on this driving "
            "domain (AUPRC = 0.22).",
            "Insight: scale alone does not guarantee better anomaly "
            "detection — training data and prompt fit matter more.",
        ],
        FIG_DIR / "fig01_binary_pr_curves.pdf",
        "Precision-Recall curves for all 15 models on the binary task.",
    ),

    # --- Slide 4: Multiclass classification ---
    (
        "Multiclass Anomaly Classification",
        "Per-class F1 across all 9 VLMs on 5,000 anomalous images",
        [
            "Now the harder task: given an anomalous image, "
            "classify the hazard into 1 of 11 classes.",
            "InternVL3-8B and Qwen2.5-VL-7B dominate, with strong "
            "F1 on common classes (animal_on_road, vehicle_incident).",
            "Smaller VLMs collapse on rare classes — InternVL3-1B and "
            "LLaVA-1.6-13B get F1 = 0 on most classes other than "
            "animals and oversized vehicles.",
            "Even the strongest models struggle on visually-similar "
            "classes (debris vs strange_object, infrastructure_failure).",
        ],
        FIG_DIR / "fig03_multiclass_f1_heatmap.pdf",
        "Per-class F1 heatmap (red = worst, green = best).",
    ),

    # --- Slide 5: Class difficulty ---
    (
        "Which Anomaly Classes Are Hardest?",
        "Mean F1 across all models per class — reveals systematic gaps",
        [
            "Animal-on-Road and Oversized Vehicle are the easiest "
            "classes (mean F1 > 0.4) — they have distinctive visual "
            "signatures and large training-data exposure.",
            "Strange Object on Road, Adverse Lighting, and "
            "Multi-Hazard are the hardest — even averaged across "
            "9 models the mean F1 is below 0.05.",
            "These hard classes are also the most safety-critical: "
            "current VLMs would fail to flag them in production.",
            "Implication: future work needs targeted data and "
            "training for rare-but-critical anomaly types.",
        ],
        APP_FIG / "appA04_hardest_classes.pdf",
        "Mean per-class F1 ± std across 9 VLMs.",
    ),

    # --- Slide 6: Calibration overview ---
    (
        "Confidence Calibration",
        "Do model-reported confidence scores reflect actual accuracy?",
        [
            "We measure Expected Calibration Error (ECE) — how well "
            "self-reported confidence predicts correctness. Lower is better.",
            "Qwen2.5-VL-3B is well-calibrated (ECE = 0.027) — "
            "its confidence numbers are actually trustworthy.",
            "InternVL3-1B is severely overconfident "
            "(reports 0.95, accuracy 0.33). Critical safety concern "
            "for AV deployment.",
            "Qwen2.5-VL-7B and LLaMA-3.2-11B are *under*-confident: "
            "they're often correct but report low confidence — "
            "their outputs would be ignored unnecessarily.",
        ],
        FIG_DIR / "fig05_calibration_ece.pdf",
        "ECE (left) and over-/under-confidence (right) per model.",
    ),

    # --- Slide 7: Reliability diagrams ---
    (
        "Reliability Diagrams",
        "Visualising calibration: 'when I say 0.8, am I right 80% of the time?'",
        [
            "Each subplot shows accuracy versus self-reported "
            "confidence per VLM. The dashed diagonal = perfect "
            "calibration.",
            "Green bars below the perfect line = underconfidence "
            "(model is more accurate than it claims).",
            "Red bars above the perfect line = overconfidence "
            "(model is less accurate than it claims).",
            "Most VLMs cluster near 0.95 confidence regardless of "
            "correctness — they exhibit a 'fixed-confidence' bias "
            "that harms calibration even when accuracy is high.",
        ],
        FIG_DIR / "fig04_reliability_diagrams.pdf",
        "Reliability diagrams for all 9 VLMs (binary detection).",
    ),

    # --- Slide 8: Cross-model agreement ---
    (
        "Cross-Model Agreement",
        "How often do VLMs make the same prediction?",
        [
            "Pairwise agreement matrix on binary detection, ordered "
            "by hierarchical clustering. Higher (greener) = more "
            "agreement.",
            "InternVL3 family (1B, 2B, 8B) clusters tightly together "
            "— scaling the same architecture preserves prediction biases.",
            "LLaMA-3.2-11B is an outlier — agrees least with everyone "
            "else, suggesting genuinely different inductive biases.",
            "Across all 15K images, no image is consistently misclassified "
            "by all 9 models — disagreement is widespread, which makes "
            "ensembling potentially valuable.",
        ],
        FIG_DIR / "fig06_agreement_matrix.pdf",
        "Pairwise binary-prediction agreement (clustering ordered).",
    ),

    # --- Slide 9: Ensemble ceiling ---
    (
        "Ensemble Upper Bound",
        "Does combining models help? — top-k majority vote",
        [
            "Adding the top-2 model (Qwen2.5-VL-3B + LLaMA-3.2-11B) "
            "boosts balanced accuracy beyond the best single model.",
            "Returns diminish quickly: by k=4 the curve plateaus, "
            "meaning the next models add redundant errors rather "
            "than complementary correctness.",
            "Best single model (Qwen2.5-VL-3B): 0.91 balanced acc. "
            "Best ensemble peak: ~0.93 — small but meaningful gain.",
            "Implication: a small, diverse ensemble of 2–3 "
            "well-chosen models is more cost-effective than running "
            "all 9.",
        ],
        FIG_DIR / "fig07_ensemble_upper_bound.pdf",
        "Top-k majority-vote ensemble accuracy as we add more models.",
    ),

    # --- Slide 10: Vision representations by class ---
    (
        "Vision-Encoder Representations: by Class",
        "PCA of the visual embedding the LLM sees, coloured by anomaly class",
        [
            "We project the projected vision-encoder output (the embedding "
            "the language model receives) of 600 anomalous images to 2D.",
            "Even in our best multiclass model (InternVL3-8B), classes "
            "do not separate cleanly in the visual embedding space.",
            "This means the *vision encoder itself* is not class-discriminative — "
            "the LLM has to do most of the disambiguation work via reasoning.",
            "Suggests a frontier for future improvement: training a "
            "vision encoder on anomaly-aware contrastive objectives "
            "could unlock big gains.",
        ],
        FIG_DIR / "fig08_vision_umap_by_class.pdf",
        "PCA of vision-encoder representations, coloured by anomaly class.",
    ),

    # --- Slide 11: Vision rep correct vs wrong ---
    (
        "Vision Reps: Correct vs Wrong Predictions",
        "Same projection, coloured by per-image binary correctness",
        [
            "For each model, green points are images it classified "
            "correctly; red are misclassifications.",
            "Strong models (InternVL3-8B, Qwen2.5-VL-3B) show "
            "correct predictions clustered together — they "
            "consistently 'see' anomalies in similar visual neighbourhoods.",
            "Weaker models (InternVL3-1B, LLaVA-1.6-13B) show errors "
            "distributed throughout the embedding space, hinting at "
            "no coherent visual notion of 'anomaly'.",
            "The visual representation is a strong predictor of "
            "correctness — supports our linear-probe finding.",
        ],
        FIG_DIR / "fig09_vision_umap_binary_grid.pdf",
        "3×3 grid: vision-rep PCA per model (green = correct, red = wrong).",
    ),

    # --- Slide 12: Layer-wise probing ---
    (
        "Layer-wise Anomaly Class Encoding",
        "Where in the vision encoder is class information stored?",
        [
            "We train a linear probe on each layer's output to "
            "predict the anomaly class — accuracy reveals where "
            "class info is encoded.",
            "All models show monotonic improvement with depth: "
            "later layers encode more class-discriminative information.",
            "Final-layer probe accuracy is much higher than zero-shot "
            "VLM accuracy on multiclass (0.50 vs 0.20–0.40), suggesting "
            "the visual features already contain the answer — the LLM "
            "fails to extract it.",
            "Strong motivation for representation-based detectors: "
            "freeze the VLM, train a small head — much better than "
            "zero-shot prompting, with 100× lower latency.",
        ],
        FIG_DIR / "fig11_linear_probe_overlay.pdf",
        "Linear-probe class-prediction accuracy per layer, all models overlaid.",
    ),

    # --- Slide 13: Model similarity (CKA) ---
    (
        "Representation Similarity Across Models",
        "Linear CKA — how similar are the learned representations?",
        [
            "Centered Kernel Alignment (CKA) measures similarity "
            "between two models' representations on the same images. "
            "1 = identical, 0 = orthogonal.",
            "InternVL3 family is internally near-identical "
            "(CKA = 0.93–0.95) despite the 8× scale difference — "
            "scale doesn't diversify representations.",
            "LLaMA-3.2-11B is an outlier (CKA ≈ 0.25–0.30 with all "
            "others) — its cross-attention vision design produces "
            "genuinely different representations.",
            "Takeaway for ensembling: combining InternVL3 variants "
            "yields little gain; combining InternVL3 with LLaMA / "
            "Qwen2.5 yields complementary errors and stronger ensembles.",
        ],
        FIG_DIR / "fig12_cka_matrix.pdf",
        "Linear CKA similarity between models on anomalous images.",
    ),
]

assert len(SLIDES) == 12, f"Expected 12 content slides, got {len(SLIDES)}"


# ──────────────────────────────────────────────────────────────────────────────
# Helper: Problem slide (left = blank space for architecture image,
# right = textual problem statement / motivation)
# ──────────────────────────────────────────────────────────────────────────────
def add_problem_slide(slide_num):
    slide = prs.slides.add_slide(BLANK)

    # Top stripe
    add_box(slide, 0, 0, SLIDE_W, Inches(0.18), fill=NAVY)

    # Title
    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.30),
                                      SLIDE_W - Inches(1.0), Inches(0.7))
    set_text(t_box.text_frame, "Problem & Motivation",
             size=28, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    # Subtitle
    s_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0),
                                     SLIDE_W - Inches(1.0), Inches(0.4))
    set_text(s_box.text_frame,
             "Why benchmarking VLMs for AV anomaly detection matters",
             size=14, color=GREY_M, align=PP_ALIGN.LEFT)

    # ── Left column: empty placeholder box for the user's architecture ─────
    LEFT_X = Inches(0.55)
    LEFT_Y = Inches(1.6)
    LEFT_W = Inches(6.0)
    LEFT_H = SLIDE_H - LEFT_Y - Inches(0.6)

    # Light dashed border so the user knows where to drop the image
    from pptx.enum.shapes import MSO_SHAPE
    placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          LEFT_X, LEFT_Y, LEFT_W, LEFT_H)
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    placeholder.line.color.rgb = GREY_L
    placeholder.line.width = Pt(1.5)
    placeholder.line.dash_style = 7  # dashed
    placeholder.shadow.inherit = False

    ph_text = slide.shapes.add_textbox(LEFT_X, LEFT_Y + LEFT_H / 2 - Inches(0.3),
                                        LEFT_W, Inches(0.6))
    set_text(ph_text.text_frame, "[ Insert architecture diagram here ]",
             size=14, color=GREY_M, align=PP_ALIGN.CENTER)

    # ── Right column: problem text ─────────────────────────────────────────
    RIGHT_X = Inches(7.0)
    RIGHT_Y = Inches(1.6)
    RIGHT_W = SLIDE_W - RIGHT_X - Inches(0.4)
    RIGHT_H = SLIDE_H - RIGHT_Y - Inches(0.6)

    text_box = slide.shapes.add_textbox(RIGHT_X, RIGHT_Y, RIGHT_W, RIGHT_H)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)

    bullets = [
        "Vision-Language Models (VLMs) are being rapidly deployed in "
        "autonomous-driving stacks for high-level scene understanding.",
        "Yet, their reliability as anomaly detectors for safety-critical "
        "deployment remains poorly characterised.",
        "Existing AV benchmarks focus on visual-question answering, "
        "trustworthiness, or robustness — but not on the specific task of "
        "detecting and characterising anomalous scenes.",
        "Open question: how do we know whether a VLM is safe to put in the "
        "perception loop of a self-driving car?",
        "Goal: build a rigorous, multi-task benchmark that exposes the "
        "deployment-relevant strengths and pathologies of state-of-the-art VLMs.",
    ]

    first = True
    for b in bullets:
        if first:
            set_text(tf, "•  " + b, size=14, color=GREY_D, line_spacing=1.35)
            first = False
        else:
            add_paragraph(tf, b, size=14, color=GREY_D, bullet=True,
                          space_before=8, line_spacing=1.35)

    # Footer
    footer = slide.shapes.add_textbox(SLIDE_W - Inches(1.0),
                                       SLIDE_H - Inches(0.4),
                                       Inches(0.8), Inches(0.3))
    set_text(footer.text_frame, f"{slide_num} / {TOTAL_SLIDES}",
             size=10, color=GREY_M, align=PP_ALIGN.RIGHT)


# ──────────────────────────────────────────────────────────────────────────────
# Helper: Text-only slide (for Conclusion and Future Work)
# Two-column layout when many bullets, single-column otherwise
# ──────────────────────────────────────────────────────────────────────────────
def add_text_slide(slide_num, title, subtitle, bullets, two_col=False):
    slide = prs.slides.add_slide(BLANK)

    # Top stripe
    add_box(slide, 0, 0, SLIDE_W, Inches(0.18), fill=NAVY)

    # Title
    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.30),
                                      SLIDE_W - Inches(1.0), Inches(0.7))
    set_text(t_box.text_frame, title, size=30, bold=True, color=NAVY,
             align=PP_ALIGN.LEFT)

    # Subtitle
    if subtitle:
        s_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.05),
                                         SLIDE_W - Inches(1.0), Inches(0.4))
        set_text(s_box.text_frame, subtitle, size=14, color=GREY_M,
                 align=PP_ALIGN.LEFT)

    BODY_TOP = Inches(1.7) if subtitle else Inches(1.4)
    BODY_H   = SLIDE_H - BODY_TOP - Inches(0.6)

    if two_col and len(bullets) >= 4:
        # Split bullets between two columns
        half = (len(bullets) + 1) // 2
        cols = [bullets[:half], bullets[half:]]
        col_w = Inches(5.9)
        col_xs = [Inches(0.55), Inches(7.0)]
        for col_x, items in zip(col_xs, cols):
            box = slide.shapes.add_textbox(col_x, BODY_TOP, col_w, BODY_H)
            tf = box.text_frame
            tf.word_wrap = True
            first = True
            for b in items:
                if first:
                    set_text(tf, "•  " + b, size=15, color=GREY_D,
                             line_spacing=1.4)
                    first = False
                else:
                    add_paragraph(tf, b, size=15, color=GREY_D, bullet=True,
                                  space_before=10, line_spacing=1.4)
    else:
        box = slide.shapes.add_textbox(Inches(0.55), BODY_TOP,
                                        SLIDE_W - Inches(1.1), BODY_H)
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for b in bullets:
            if first:
                set_text(tf, "•  " + b, size=16, color=GREY_D,
                         line_spacing=1.4)
                first = False
            else:
                add_paragraph(tf, b, size=16, color=GREY_D, bullet=True,
                              space_before=12, line_spacing=1.4)

    # Footer
    footer = slide.shapes.add_textbox(SLIDE_W - Inches(1.0),
                                       SLIDE_H - Inches(0.4),
                                       Inches(0.8), Inches(0.3))
    set_text(footer.text_frame, f"{slide_num} / {TOTAL_SLIDES}",
             size=10, color=GREY_M, align=PP_ALIGN.RIGHT)


# ──────────────────────────────────────────────────────────────────────────────
# Build slides in the right order
# ──────────────────────────────────────────────────────────────────────────────

# Slide 2: Problem
add_problem_slide(2)
print("  Slide  2: Problem & Motivation")

# Slides 3-14: 12 content slides
for i, (title, subtitle, bullets, fig_path, caption) in enumerate(SLIDES, start=3):
    add_content_slide(i, title, subtitle, bullets, fig_path, caption)
    print(f"  Slide {i:2d}: {title}")

# Slide 15: Conclusion
conclusion_bullets = [
    "We introduced → DASHA-15K: a 15{,}000-image dashcam benchmark with "
    "10K normal + 5K anomalous images across 11 hazard classes, with "
    "rich human-curated and cross-LLM-validated annotations.",
    "We evaluated 16 models — 9 open VLMs (1B-13B) plus 6 visual-similarity "
    "baselines (CLIP, SigLIP) and a reconstruction autoencoder — across "
    "binary detection, multiclass classification, and free-form description.",
    "Three deployment-relevant pathologies emerge:",
    "    • Scale ≠ Better: a 3B VLM (Qwen2.5-VL-3B) is the best binary "
    "detector, beating models 4× larger.",
    "    • Confidence is unreliable: 7 of 9 VLMs are mis-calibrated, with "
    "severe over- and under-confidence patterns that pose direct safety hazards.",
    "    • Within-family representations are collinear: scaling alone "
    "(InternVL3 1B → 8B; CKA > 0.93) does not diversify what the model sees.",
    "Linear probes on frozen vision-encoder features beat zero-shot prompting "
    "on multiclass classification — a high-impact follow-up for low-latency, "
    "deployable AV anomaly detection.",
    "All data, code, predictions, and 135K saved hidden-state representations "
    "released under CC-BY 4.0.",
]
add_text_slide(15, "Conclusion",
               "Key takeaways from the DASHA-15K benchmark study",
               conclusion_bullets, two_col=False)
print("  Slide 15: Conclusion")

# Slide 16: Future Work
future_bullets = [
    "Train a small classification head on frozen vision-encoder features — "
    "linear probes already suggest 100× lower latency than zero-shot prompting "
    "and competitive accuracy.",
    "Extend the benchmark to temporal/video anomalies (e.g., gradual lane drift, "
    "approaching hazards) with multi-frame inputs.",
    "Evaluate proprietary VLMs (GPT-4V, Gemini, Claude) on the full benchmark "
    "and compare against the open-source zoo.",
    "Develop post-hoc calibration techniques tailored to the over/under-confidence "
    "patterns we identified — a prerequisite for safe deployment.",
    "Cross-domain transfer: how do these models generalise from urban driving "
    "to highway, off-road, or adverse-weather settings?",
    "Fine-tuning on the training portion of DASHA-15K — our zero-shot numbers "
    "are likely a lower bound on what these models can achieve.",
    "Architecture-diverse ensembles: pair Qwen2.5-VL-3B with the "
    "cross-attention-based LLaMA-3.2-11B for the strongest binary detector "
    "we observed (top-2 ensemble bal. acc. 0.93).",
]
add_text_slide(16, "Future Work",
               "Where to go next, building on DASHA-15K",
               future_bullets, two_col=True)
print("  Slide 16: Future Work")

# Save
out_path = PRES_DIR / "AV_VLM_Benchmark.pptx"
prs.save(out_path)

print(f"\nPresentation saved → {out_path}")
print(f"Slides total       : {len(prs.slides)}")
